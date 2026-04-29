
import logging
import base64
import io
from pathlib import Path
from typing import Optional, List, Dict, Tuple

logger = logging.getLogger(__name__)


async def extract_text_and_images_from_pdf(file_path: str) -> Tuple[str, List[Dict]]:
    try:
        import pdfplumber
        from PIL import Image
        import fitz

        text_parts = []
        images = []

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        try:
            pdf_doc = fitz.open(file_path)
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                image_list = page.get_images(full=True)

                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = pdf_doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]

                    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    images.append({
                        "page": page_num + 1,
                        "index": img_index,
                        "format": image_ext,
                        "data": image_base64,
                        "description": f"Image from page {page_num + 1}"
                    })

            pdf_doc.close()
        except Exception as img_error:
            logger.warning(f"Failed to extract images from PDF (text extraction succeeded): {str(img_error)}")

        extracted_text = "\n\n".join(text_parts)
        logger.info(f"Successfully extracted {len(extracted_text)} characters and {len(images)} images from PDF: {file_path}")
        return extracted_text, images

    except Exception as e:
        logger.error(f"Failed to extract text from PDF {file_path}: {str(e)}")
        raise Exception(f"PDF text extraction failed: {str(e)}")


async def extract_text_from_pdf(file_path: str) -> str:
    text, _ = await extract_text_and_images_from_pdf(file_path)
    return text


async def extract_text_and_images_from_pptx(file_path: str) -> Tuple[str, List[Dict]]:
    try:
        from pptx import Presentation
        import zipfile
        import base64

        prs = Presentation(file_path)
        text_parts = []
        images = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells if cell.text)
                        if row_text:
                            slide_text_parts.append(row_text)

                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext

                        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                        images.append({
                            "slide": slide_num,
                            "format": image_ext,
                            "data": image_base64,
                            "description": f"Image from slide {slide_num}"
                        })
                    except Exception as img_error:
                        logger.debug(f"Failed to extract image from shape: {str(img_error)}")

            if slide_text_parts:
                slide_text = "\n".join(slide_text_parts)
                text_parts.append(slide_text)

        extracted_text = "\n\n".join(text_parts)
        logger.info(f"Successfully extracted {len(extracted_text)} characters and {len(images)} images from PPTX: {file_path}")
        return extracted_text, images

    except Exception as e:
        logger.error(f"Failed to extract text from PPTX {file_path}: {str(e)}")
        raise Exception(f"PPTX text extraction failed: {str(e)}")


async def extract_text_from_pptx(file_path: str) -> str:
    text, _ = await extract_text_and_images_from_pptx(file_path)
    return text


async def extract_text_and_images_from_docx(file_path: str) -> Tuple[str, List[Dict]]:
    try:
        from docx import Document
        import zipfile
        from PIL import Image
        import base64

        doc = Document(file_path)
        text_parts = []
        images = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        try:
            docx_zip = zipfile.ZipFile(file_path, 'r')
            image_files = [f for f in docx_zip.namelist() if f.startswith('word/media/') and f.split('.')[-1].lower() in ['png', 'jpg', 'jpeg', 'gif', 'bmp']]

            for img_index, img_path in enumerate(image_files):
                image_data = docx_zip.read(img_path)
                image_ext = img_path.split('.')[-1].lower()

                image_base64 = base64.b64encode(image_data).decode('utf-8')
                images.append({
                    "index": img_index,
                    "format": image_ext,
                    "data": image_base64,
                    "description": f"Image {img_index + 1} from document"
                })

            docx_zip.close()
        except Exception as img_error:
            logger.warning(f"Failed to extract images from DOCX (text extraction succeeded): {str(img_error)}")

        extracted_text = "\n".join(text_parts)

        if not extracted_text or not extracted_text.strip():
            logger.warning(f"DOCX file {file_path} appears to be empty or contains no extractable text")
            try:
                import zipfile
                with zipfile.ZipFile(file_path, 'r') as z:
                    files = z.namelist()
                    logger.info(f"DOCX contains {len(files)} files")
                    if 'word/document.xml' not in files:
                        logger.error(f"DOCX file {file_path} is missing word/document.xml - may be corrupted")
                    else:
                        try:
                            doc_xml = z.read('word/document.xml')
                            logger.info(f"document.xml size: {len(doc_xml)} bytes")
                            if b'<w:t>' in doc_xml:
                                logger.warning(f"document.xml contains <w:t> tags but no text was extracted - possible issue with python-docx library")
                        except Exception as xml_error:
                            logger.error(f"Could not read document.xml: {str(xml_error)}")
            except Exception as check_error:
                logger.error(f"Failed to validate DOCX file structure: {str(check_error)}")
        else:
            logger.info(f"Successfully extracted {len(extracted_text)} characters from DOCX")

        logger.info(f"Extraction complete: {len(extracted_text)} characters and {len(images)} images from DOCX: {file_path}")
        return extracted_text, images

    except Exception as e:
        logger.error(f"Failed to extract text from DOCX {file_path}: {str(e)}", exc_info=True)
        return "", []


async def extract_text_from_docx(file_path: str) -> str:
    text, _ = await extract_text_and_images_from_docx(file_path)
    return text


async def extract_text_from_txt(file_path: str) -> Tuple[str, List[Dict]]:
    try:
        encodings = ['utf-8', 'utf-8-sig', 'cp1251', 'latin-1', 'iso-8859-1']
        text_content = None

        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text_content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if text_content is None:
            with open(file_path, 'rb') as f:
                text_content = f.read().decode('utf-8', errors='replace')

        logger.info(f"Successfully extracted {len(text_content)} characters from TXT: {file_path}")
        return text_content, []

    except Exception as e:
        logger.error(f"Failed to extract text from TXT {file_path}: {str(e)}")
        raise Exception(f"TXT text extraction failed: {str(e)}")


async def process_lesson_material(
    file_path: str,
    material_type: str,
    extract_images: bool = True
) -> Tuple[Optional[str], List[Dict]]:
    material_type = material_type.upper()

    if material_type not in ["PDF", "PPTX", "DOCX", "TXT"]:
        logger.info(f"Skipping text extraction for unsupported type: {material_type}")
        return None, []

    if not Path(file_path).exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    if material_type == "PDF":
        if extract_images:
            return await extract_text_and_images_from_pdf(file_path)
        else:
            text = await extract_text_from_pdf(file_path)
            return text, []
    elif material_type == "PPTX":
        if extract_images:
            return await extract_text_and_images_from_pptx(file_path)
        else:
            text = await extract_text_from_pptx(file_path)
            return text, []
    elif material_type == "DOCX":
        if extract_images:
            return await extract_text_and_images_from_docx(file_path)
        else:
            text = await extract_text_from_docx(file_path)
            return text, []
    elif material_type == "TXT":
        return await extract_text_from_txt(file_path)

    return None, []


async def process_multiple_materials(materials: list, db_session=None) -> str:
    from app.models.lesson_material import MaterialType
    from app.core.storage import storage_client

    text_parts = []

    for material in materials:
        if not material.use_for_ai_generation:
            logger.debug(f"Skipping material {material.id} - use_for_ai_generation=False")
            continue

        logger.info(f"Processing material {material.id}: type={material.type.value}, title={material.title}")

        try:
            if material.type == MaterialType.TEXT:
                if material.content and material.content.strip():
                    text_parts.append(f"=== {material.title} ===\n{material.content}")
                    logger.info(f"Added text content from material {material.id}: {len(material.content)} characters")
                else:
                    logger.warning(f"Material {material.id} (TEXT) has no content")
                continue

            if material.extracted_text and material.extracted_text.strip():
                material_text = f"=== {material.title} ===\n{material.extracted_text}"

                if material.extracted_images:
                    try:
                        import json
                        images_data = json.loads(material.extracted_images)
                        if images_data:
                            material_text += f"\n\n[В документе содержится {len(images_data)} изображений, которые иллюстрируют материал]"
                            for img in images_data:
                                if "description" in img:
                                    material_text += f"\n- {img['description']}"
                    except Exception as e:
                        logger.warning(f"Failed to parse images for material {material.id}: {str(e)}")

                text_parts.append(material_text)
                logger.info(f"Used extracted_text from material {material.id}: {len(material.extracted_text)} characters")
                continue

            if material.type in [MaterialType.PDF, MaterialType.DOCX, MaterialType.PPTX, MaterialType.TXT]:
                if material.file_url:
                    if not material.extracted_text or not material.extracted_text.strip():
                        logger.info(f"Material {material.id} has no extracted_text, extracting from file {material.file_url}...")
                        import tempfile
                        import os

                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{material.type.value}") as tmp_file:
                            tmp_path = tmp_file.name

                            try:
                                logger.info(f"Downloading file from MinIO: {material.file_url}")
                                file_data = storage_client.download_file(material.file_url)
                                tmp_file.write(file_data)
                                tmp_file.flush()
                                logger.info(f"File downloaded, size: {len(file_data)} bytes")

                                logger.info(f"Extracting text from {material.type.value} file...")
                                extracted_text, extracted_images = await process_lesson_material(
                                    tmp_path,
                                    material.type.value,
                                    extract_images=True
                                )

                                logger.info(f"Extraction result: text_length={len(extracted_text) if extracted_text else 0}, images_count={len(extracted_images) if extracted_images else 0}")

                                if extracted_text and extracted_text.strip() and db_session:
                                    try:
                                        from sqlalchemy import update
                                        from app.models.lesson_material import LessonMaterial

                                        update_stmt = update(LessonMaterial).where(
                                            LessonMaterial.id == material.id
                                        ).values(
                                            extracted_text=extracted_text
                                        )

                                        await db_session.execute(update_stmt)
                                        await db_session.commit()
                                        logger.info(f"Saved extracted_text ({len(extracted_text)} chars) to material {material.id} in database")
                                    except Exception as save_error:
                                        logger.warning(f"Failed to save extracted_text to database for material {material.id}: {str(save_error)}")

                                material_text = f"=== {material.title} ===\n"
                                if extracted_text and extracted_text.strip():
                                    material_text += extracted_text
                                    logger.info(f"Successfully extracted {len(extracted_text)} characters from {material.type.value} file")
                                else:
                                    logger.warning(f"No text extracted from {material.type.value} file (file might be empty or corrupted)")

                                if extracted_images:
                                    material_text += f"\n\n[В документе содержится {len(extracted_images)} изображений, которые иллюстрируют материал]"
                                    for img in extracted_images:
                                        if "description" in img:
                                            material_text += f"\n- {img['description']}"

                                if material_text.strip():
                                    text_parts.append(material_text)
                                    logger.info(f"Added text from file material {material.id}: {len(extracted_text) if extracted_text else 0} characters")
                                else:
                                    logger.error(f"Material {material.id}: No text content after extraction from {material.type.value} file")
                            except Exception as e:
                                error_msg = f"Failed to extract text from file material {material.id}: {str(e)}"
                                logger.error(error_msg, exc_info=True)
                                import traceback
                                logger.error(f"Full traceback: {traceback.format_exc()}")
                            finally:
                                if os.path.exists(tmp_path):
                                    try:
                                        os.unlink(tmp_path)
                                    except:
                                        pass
                    else:
                        logger.info(f"Material {material.id} already has extracted_text ({len(material.extracted_text)} characters), using it")
                else:
                    logger.warning(f"Material {material.id} ({material.type.value}) has no file_url, cannot extract text")

        except Exception as e:
            logger.warning(f"Failed to process material {material.id}: {str(e)}")
            continue

    combined_text = "\n\n".join(text_parts)
    logger.info(f"Combined {len(materials)} materials into {len(combined_text)} characters of text")

    processed_count = len(text_parts)
    logger.info(f"Successfully processed {processed_count} out of {len(materials)} materials")

    if len(combined_text.strip()) < 100:
        logger.warning(f"Combined text is very short ({len(combined_text)} characters), may not be sufficient for AI generation")
        for material in materials:
            logger.warning(f"Material {material.id}: type={material.type.value}, "
                          f"has_content={bool(material.content)}, "
                          f"has_extracted_text={bool(material.extracted_text)}, "
                          f"extracted_text_length={len(material.extracted_text) if material.extracted_text else 0}, "
                          f"has_file_url={bool(material.file_url)}, "
                          f"file_url={material.file_url if material.file_url else 'None'}, "
                          f"use_for_ai={material.use_for_ai_generation}")

    return combined_text
